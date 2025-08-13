#!/usr/bin/env python3
"""
Extract PowerPoint slides as images using python-pptx and Pillow
"""

import sys
import json
import os
from pptx import Presentation
from PIL import Image
import io
import base64

def extract_slides_as_images(pptx_path):
    """Extract all slides from PowerPoint as images"""
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                'number': i,
                'title': '',
                'bullets': [],
                'images': [],
                'hasBackground': False
            }
            
            # Extract text content
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if shape == slide.shapes.title:
                            slide_data['title'] = text
                        else:
                            # Extract bullet points
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    slide_data['bullets'].append(para_text)
                
                # Check for images
                if hasattr(shape, 'image'):
                    slide_data['images'].append({
                        'type': 'embedded',
                        'format': shape.image.ext
                    })
            
            # Check for background
            if slide.background and slide.background.fill:
                slide_data['hasBackground'] = True
            
            slides_data.append(slide_data)
        
        return {
            'success': True,
            'slides': slides_data,
            'total': len(slides_data)
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'slides': []
        }

if __name__ == '__main__':
    if len(sys.argv) != 2:
        print(json.dumps({'error': 'Usage: python extract_slides_as_images.py <pptx_file>'}))
        sys.exit(1)
    
    result = extract_slides_as_images(sys.argv[1])
    print(json.dumps(result))