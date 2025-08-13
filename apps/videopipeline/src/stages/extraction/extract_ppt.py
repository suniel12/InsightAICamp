#!/usr/bin/env python3
"""
PowerPoint extraction script using python-pptx
"""

import sys
import json
from pptx import Presentation

def extract_text_from_shape(shape):
    """Extract text from a PowerPoint shape - optimized version"""
    # Best practice: Check for text attribute to avoid exceptions
    if not hasattr(shape, "text"):
        return ""
    
    # Store text in variable to avoid repeated access (performance optimization)
    shape_text = shape.text
    return shape_text if shape_text else ""

def extract_slide_content(slide):
    """Extract all content from a slide with performance optimizations"""
    slide_data = {
        "title": "",
        "bullets": [],
        "notes": "",
        "images": [],
        "tables": [],
        "charts": []
    }
    
    # Extract title - store in variable to avoid repeated access
    if slide.shapes.title:
        title_text = slide.shapes.title.text
        slide_data["title"] = title_text if title_text else ""
    
    # Extract text from all shapes - optimized iteration
    for shape in slide.shapes:
        # Handle text frames with proper run extraction
        if shape.has_text_frame:
            # Best practice: Extract text from runs for complete formatting
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                # Join all runs in paragraph to preserve formatting boundaries
                para_text = ''.join(run.text for run in paragraph.runs)
                if para_text.strip():
                    text_parts.append(para_text.strip())
            
            # Only add non-empty, non-title text
            if text_parts and '\n'.join(text_parts) != slide_data["title"]:
                slide_data["bullets"].extend(text_parts)
        
        # Check for images with error handling
        try:
            if shape.shape_type == 13:  # Picture
                slide_data["images"].append({
                    "type": "embedded",
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                })
        except Exception as e:
            # Log but don't fail on individual shape errors
            print(f"Warning: Could not extract image: {e}", file=sys.stderr)
        
        # Extract tables if present
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            slide_data["tables"].append(table_data)
    
    # Extract speaker notes with error handling
    try:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            slide_data["notes"] = notes_text if notes_text else ""
    except Exception as e:
        print(f"Warning: Could not extract notes: {e}", file=sys.stderr)
    
    return slide_data

def extract_presentation(pptx_path):
    """Extract all content from a PowerPoint presentation"""
    try:
        presentation = Presentation(pptx_path)
        
        extracted_data = {
            "title": "",
            "author": "",
            "slides": []
        }
        
        # Extract metadata
        if presentation.core_properties.title:
            extracted_data["title"] = presentation.core_properties.title
        if presentation.core_properties.author:
            extracted_data["author"] = presentation.core_properties.author
        
        # Extract slides
        for slide in presentation.slides:
            slide_content = extract_slide_content(slide)
            extracted_data["slides"].append(slide_content)
        
        # If no title found in metadata, use first slide title
        if not extracted_data["title"] and extracted_data["slides"]:
            extracted_data["title"] = extracted_data["slides"][0]["title"]
        
        return extracted_data
    
    except Exception as e:
        return {
            "error": str(e),
            "slides": []
        }

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Usage: python extract_ppt.py <pptx_file>"}))
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    result = extract_presentation(pptx_file)
    print(json.dumps(result, indent=2))